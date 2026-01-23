import io
import os
import re
import pymupdf
import pymupdf4llm
import pytesseract
from PIL import Image
from weasyprint import HTML
from docx import Document
from pptx import Presentation

# Ensure pytesseract path is set if needed (optional for some envs)
# pytesseract.pytesseract.tesseract_cmd = r'/usr/bin/tesseract' 

def convert_to_pdf(input_path, output_pdf_path):
    """
    Converts .docx, .pptx, .txt, .png, .jpg to a unified PDF format.
    Uses pure Python libraries (WeasyPrint, python-docx, python-pptx, Pillow).
    """
    ext = os.path.splitext(input_path)[1].lower()

    try:
        if ext == '.pdf':
            # Just copy or simple pass-through. 
            # We use pymupdf to save it to ensure it's a valid clean PDF.
            doc = pymupdf.open(input_path)
            doc.save(output_pdf_path)
            doc.close()

        elif ext in ['.png', '.jpg', '.jpeg']:
            image = Image.open(input_path)
            # Convert to RGB to avoid alpha channel issues in PDF
            if image.mode in ("RGBA", "P"):
                image = image.convert("RGB")
            image.save(output_pdf_path, "PDF", resolution=100.0)

        elif ext == '.txt':
            with open(input_path, 'r', encoding='utf-8') as f:
                text = f.read()
            # Wrap in simple HTML
            html_content = f"<html><body><pre>{text}</pre></body></html>"
            HTML(string=html_content).write_pdf(output_pdf_path)

        elif ext == '.docx':
            doc = Document(input_path)
            html_parts = ["<html><body>"]
            for para in doc.paragraphs:
                if para.text.strip():
                    html_parts.append(f"<p>{para.text}</p>")
            # Handle simple tables (optional but good for assignments)
            for table in doc.tables:
                html_parts.append("<table border='1'>")
                for row in table.rows:
                    html_parts.append("<tr>")
                    for cell in row.cells:
                        html_parts.append(f"<td>{cell.text}</td>")
                    html_parts.append("</tr>")
                html_parts.append("</table>")
            
            html_parts.append("</body></html>")
            HTML(string="".join(html_parts)).write_pdf(output_pdf_path)

        elif ext == '.pptx':
            prs = Presentation(input_path)
            html_parts = ["<html><body>"]
            for slide in prs.slides:
                html_parts.append("<div style='page-break-after: always; border: 1px solid #ccc; padding: 20px;'>")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Simple rendering of text
                        html_parts.append(f"<p>{shape.text}</p>")
                    # Note: Extracting images from PPTX to HTML is complex; 
                    # we focus on text here. If user needs PPTX images, 
                    # we'd need to extract them to temp dir and link them.
                html_parts.append("</div>")
            html_parts.append("</body></html>")
            HTML(string="".join(html_parts)).write_pdf(output_pdf_path)
        
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    except Exception as e:
        raise RuntimeError(f"Conversion failed for {ext}: {str(e)}")

def extract_with_diagram_text(pdf_path):
    """
    Extracts Markdown text from PDF, identifying "omitted" images 
    and running OCR on them to insert text back into the Markdown.
    """
    doc = pymupdf.open(pdf_path)
    
    try:
        # 1. Get the initial Markdown (this will have the "omitted" tags)
        # Pass the doc object to avoid opening the file twice
        md_text = pymupdf4llm.to_markdown(doc)
        
        # 2. Find every "picture omitted" tag in the markdown
        # Example tag: ==> picture [409 x 197] intentionally omitted <==
        # Adjust pattern to match pymupdf4llm output variations if any
        pattern = r"==> picture \[\d+ x \d+\] intentionally omitted <=="
        placeholders = re.findall(pattern, md_text)
        
        # 3. Iterate through pages and images to find the content
        img_counter = 0
        
        # Check if we need fallback mode (No placeholders found, but images might exist)
        use_fallback_ocr = len(placeholders) == 0

        # Loop over pages
        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images(full=True)
            
            for img in images:
                # If NOT in fallback mode, stop if we matched all placeholders
                if not use_fallback_ocr and img_counter >= len(placeholders): 
                    break
                
                # Extract and OCR the image
                xref = img[0]
                try:
                    pix = pymupdf.Pixmap(doc, xref)
                    # Handle CMYK or other non-RGB colorspaces if necessary
                    if pix.n - pix.alpha < 3: 
                        # e.g. grayscale, convert to RGB for consistency
                        pass 
                    
                    # If pix is CMYK, convert to RGB
                    if pix.n >= 4:
                         pix = pymupdf.Pixmap(pymupdf.csRGB, pix)

                    img_bytes = pix.tobytes()
                    img_data = Image.open(io.BytesIO(img_bytes))
                    
                    # Use Tesseract to "read" the diagram (AND, OR, C1, etc.)
                    ocr_result = pytesseract.image_to_string(img_data).strip()
                    
                    # Clean up and format the result for the AI
                    if ocr_result:
                        formatted_ocr = f"\n> **[Extracted from Image]:** {ocr_result.replace('\n', ' ')}\n"
                    else:
                        formatted_ocr = "\n> *[Image contains no readable text]*\n"
                    
                    if use_fallback_ocr:
                        # Append to end of text
                        md_text += formatted_ocr
                    else:
                        # Replace the placeholder
                        md_text = md_text.replace(placeholders[img_counter], formatted_ocr, 1)
                    
                except Exception as e:
                    print(f"OCR Error on image {img_counter}: {e}")
                    # Fallback: remove placeholder if present
                    if not use_fallback_ocr:
                        md_text = md_text.replace(placeholders[img_counter], "> *[Error processing image]*", 1)

                finally:
                    img_counter += 1
                    pix = None # Clear memory

        return md_text
    
    finally:
        if doc:
            doc.close()